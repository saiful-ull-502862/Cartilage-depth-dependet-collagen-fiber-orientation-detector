
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_algorithm_slide(output_file):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    BLUE = RGBColor(0, 51, 153)
    GREY = RGBColor(128, 128, 128)
    BLACK = RGBColor(0, 0, 0)
    
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    p = tbox.text_frame.add_paragraph()
    p.text = "Scientific Architecture: Image Processing Algorithm"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # --- Column 1: HSV Conversion ---
    l1 = Inches(0.5)
    top = Inches(1.5)
    w = Inches(4.0)
    h = Inches(5.0)
    
    # Panel Title
    t1 = slide.shapes.add_textbox(l1, top, w, Inches(0.5))
    t1.text_frame.text = "1. HSV Color Space Transformation"
    t1.text_frame.paragraphs[0].font.bold = True
    
    # Text Description
    desc1 = slide.shapes.add_textbox(l1, top + Inches(3.5), w, Inches(2.0))
    tf1 = desc1.text_frame
    tf1.word_wrap = True
    p = tf1.add_paragraph()
    p.text = "• BGR pixels are mapped to HSV cylinders."
    p.font.size = Pt(14)
    p = tf1.add_paragraph()
    p.text = "• Hue (θ) represents the dominant spectral wavelength, directly correlated to collagen birefringence."
    p.font.size = Pt(14)

    # --- Column 2: Noise Masking ---
    l2 = Inches(4.8)
    
    t2 = slide.shapes.add_textbox(l2, top, w, Inches(0.5))
    t2.text_frame.text = "2. Dual-Threshold Noise Masking"
    t2.text_frame.paragraphs[0].font.bold = True
    
    desc2 = slide.shapes.add_textbox(l2, top + Inches(3.5), w, Inches(2.0))
    tf2 = desc2.text_frame
    p = tf2.add_paragraph()
    p.text = "• Filters out background and non-tissue artifacts."
    p.font.size = Pt(14)
    p = tf2.add_paragraph()
    p.text = "• Condition: Exclude if Intensity (V) < 10 or Saturation (S) < 10."
    p.font.size = Pt(14)

    # --- Column 3: Calibration ---
    l3 = Inches(9.1)
    
    t3 = slide.shapes.add_textbox(l3, top, w, Inches(0.5))
    t3.text_frame.text = "3. Auto-Calibration (Anchors)"
    t3.text_frame.paragraphs[0].font.bold = True
    
    desc3 = slide.shapes.add_textbox(l3, top + Inches(3.5), w, Inches(2.0))
    tf3 = desc3.text_frame
    p = tf3.add_paragraph()
    p.text = "• Histogram Analysis identifies dominant modes."
    p.font.size = Pt(14)
    p = tf3.add_paragraph()
    p.text = "• H₀ = Median(SZ Region)"
    p.font.size = Pt(14)
    p = tf3.add_paragraph()
    p.text = "• H₉₀ = Median(DZ Region)"
    p.font.size = Pt(14)
    
    # --- Insert Generated Plots if they exist ---
    # We will try to insert the pngs we just made
    
    # Plot 1: Transfer Function (Put in Calib column or Split?)
    # Let's put Histogram in Col 3
    if os.path.exists('scientific_anchor_histogram.png'):
        slide.shapes.add_picture('scientific_anchor_histogram.png', l3, top + Inches(0.5), width=w)
        
    # Plot 2: Transfer Function (Put in Col 1 as it explains Hue mapping?)
    if os.path.exists('scientific_transfer_function.png'):
        slide.shapes.add_picture('scientific_transfer_function.png', l1, top + Inches(0.5), width=w)
        
    # For col 2, we don't have a plot, so we leave it for the user to insert the schematic image,
    # OR we insert a placeholder
    
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l2, top + Inches(0.5), w, Inches(2.8))
    ph.text_frame.text = "[Insert Noise Masking Diagram Here]"
    
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    try:
        create_algorithm_slide("Algorithm_Scientific_Detail.pptx")
    except Exception as e:
        print(f"Error: {e}")
