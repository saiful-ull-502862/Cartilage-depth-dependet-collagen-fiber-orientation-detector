
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

def create_interface_workflow(output_file):
    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Specific Colors matching the request
    GREEN = RGBColor(0, 153, 0)
    RED = RGBColor(204, 0, 0)
    BLUE = RGBColor(0, 76, 153)
    YELLOW = RGBColor(255, 192, 0)
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)

    # --- Slide 1: Interface Workflow ---
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = title_shape.text_frame
    p = tf.add_paragraph()
    p.text = "PLM Analysis Interface Workflow"
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = BLACK

    # Helper for boxes
    def add_process_box(left, top, width, height, color, title, steps, icon_shape=None):
        # Frame
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rect.fill.background()
        rect.line.color.rgb = color
        rect.line.width = Pt(3)
        
        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.5))
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        header.text_frame.text = title
        header.text_frame.paragraphs[0].font.bold = True
        header.text_frame.paragraphs[0].font.color.rgb = WHITE
        header.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Body
        tb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.6), width - Inches(0.2), height - Inches(0.7))
        tf = tb.text_frame
        for step in steps:
            p = tf.add_paragraph()
            p.text = step
            p.font.size = Pt(11)
            p.space_after = Pt(8)

    # === LEFT COLUMN: Input & Split ===
    # 1. Start: Image
    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(2.0), Inches(0.8)).text_frame.text = "1. PLM Image Upload\n(BGR Format)"
    
    # Arrow
    slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.4), Inches(2.35), Inches(0.2), Inches(0.4))
    
    # 2. ROI
    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.8), Inches(2.0), Inches(0.8)).text_frame.text = "2. ROI Selection\n(Zoom & Crop)"
    
    # Split Paths using Elbow Connectors (simulated with lines) or just arrows
    # Green Path (Left-ish)
    gr_network = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, Inches(1.5), Inches(3.6), Inches(3.2), Inches(2.5)) # Points to Ref Box
    gr_network.line.color.rgb = GREEN
    gr_network.line.width = Pt(3)
    
    # Red Path (Right-ish)
    rd_network = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, Inches(1.5), Inches(3.6), Inches(3.2), Inches(5.5)) # Points to Batch Box
    rd_network.line.color.rgb = RED
    rd_network.line.width = Pt(3)

    # Label paths
    lbl1 = slide.shapes.add_textbox(Inches(1.8), Inches(3.8), Inches(1), Inches(0.5))
    lbl1.text_frame.text = "Control\n(Reference)"
    lbl1.text_frame.paragraphs[0].font.color.rgb = GREEN
    
    lbl2 = slide.shapes.add_textbox(Inches(1.8), Inches(4.5), Inches(1), Inches(0.5))
    lbl2.text_frame.text = "Degraded\n(Batch)"
    lbl2.text_frame.paragraphs[0].font.color.rgb = RED


    # === MIDDLE COLUMN: Interface Processing Logic ===
    
    # Box 1: Reference Logic (Green)
    add_process_box(Inches(3.2), Inches(1.5), Inches(4.5), Inches(2.8), GREEN, "3A. Control Reference Analysis", [
        "• HSV Conversion (Hue: 0-179)",
        "• Noise Masking: Exclude V<10, S<10",
        "• Zone Separation (Normalized):",
        "   - SZ (0-10%), MZ (10-40%), DZ (40-100%)",
        "• Anchor Extraction:",
        "   - H₀ = Median Hue (SZ)",
        "   - H₉₀ = Median Hue (DZ)",
        "• Output: Calibration Map [H₀, H₉₀]"
    ])
    
    # Connector "Map" from Green to Red
    conn = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(5.45), Inches(4.3), Inches(0.3), Inches(0.7))
    conn.text_frame.text = "Params"
    conn.text_frame.paragraphs[0].font.size = Pt(8)

    # Box 2: Batch Logic (Red)
    add_process_box(Inches(3.2), Inches(5.0), Inches(4.5), Inches(2.2), RED, "3B. Degraded Batch Processing", [
        "• Apply Map [H₀, H₉₀] from Reference",
        "• Standardized Zonal Segmentation",
        "• Angle Calculation (θ):",
        "   θ = [(H - H₀) / (H₉₀ - H₀)] × 90",
        "• Circular Mean Smoothing"
    ])

    # === RIGHT COLUMN: Outputs ===
    
    # Arrow Right
    slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.7), Inches(6.0), Inches(0.5), Inches(0.3))

    # Box 3: Visuals (Yellow/Orange)
    add_process_box(Inches(8.5), Inches(1.5), Inches(4.5), Inches(2.5), YELLOW, "4. Visual Output", [
        "• Annotated Images",
        "• Zonal Boundaries Overlay",
        "• Validated Pixel Masks"
    ])
    
    # Box 4: Quantitative (Blue)
    add_process_box(Inches(8.5), Inches(4.5), Inches(4.5), Inches(2.7), BLUE, "5. Quantitative Data", [
        "• Depth Profile Graph (0-100%):",
        "   - Fiber Angle vs Normalized Depth",
        "• Comparisons:",
        "   - Healthy Curve vs Degraded Curve",
        "• Excel Export:",
        "   - Raw Binned Data & Statistics"
    ])

    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    try:
        create_interface_workflow("Interface_Process_Workflow.pptx")
    except ImportError:
        print("Error: python-pptx library is missing.")
    except Exception as e:
        print(f"An error occurred: {e}")
