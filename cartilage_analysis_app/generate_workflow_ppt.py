
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(output_file):
    prs = Presentation()

    # defined colors
    BLUE = RGBColor(0, 51, 153)
    BLACK = RGBColor(0, 0, 0)

    # Helper to add a slide with title and content
    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = BLUE
        title.text_frame.paragraphs[0].font.bold = True
        
        # Content
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # clear existing
        
        for item in content_text_list:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(22)
            p.level = 0
            p.space_after = Pt(12)

    # --- Slide 1: Title ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    
    title.text = "Methodology: Collagen Fibril Orientation Analysis"
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Image Processing & Quantitative Comparative Analysis\nBased on Polarized Light Microscopy (PLM)"
    
    # --- Slide 2: Workflow Overview ---
    add_slide("Workflow Algorithm", [
        "1. Image Acquisition & Preprocessing (HSV Conversion)",
        "2. ROI Selection & Noise Masking",
        "3. Standardized Zone Segmentation (SZ, MZ, DZ)",
        "4. Reference-Based Color Calibration",
        "5. Mathematical Angle Calculation",
        "6. Comparative Analysis (Healthy vs. Degraded)"
    ])
    
    # --- Slide 3: Step 1 - Processing ---
    add_slide("1. Image Processing & Masking", [
        "• Input: PLM Image (BGR format).",
        "• Conversion: Transformed to HSV Color Space.",
        "    - Hue (0-179): Encodes fibril orientation.",
        "    - Saturation (S) & Value (V): Used for filtering.",
        "• Noise Reduction: Dual-threshold mask applied.",
        "    - Excluded pixels where S < 10 or V < 10.",
        "    - Ensures analysis only on birefringent tissue."
    ])
    
    # --- Slide 4: Step 2 - Zonal Segmentation ---
    add_slide("2. Standardized Zone Segmentation", [
        "• Zones defined by normalized tissue thickness:",
        "    - Superficial Zone (SZ): 0% – 10%",
        "    - Middle Zone (MZ): 10% – 40%",
        "    - Deep Zone (DZ): 40% – 100%",
        "• Purpose: Ensures anatomical consistency across all samples (Femoral, Tibial, etc)."
    ])
    
    # --- Slide 5: Step 3 - Calibration ---
    add_slide("3. Reference-Based Calibration", [
        "• Reference Image: One healthy sample selected per group (e.g., Healthy Femur).",
        "• Anchors Defined:",
        "    - H₀: Median Hue of SZ (mapped to 0°).",
        "    - H₉₀: Median Hue of DZ (mapped to 90°).",
        "• Application: These H₀ and H₉₀ values are applied to all degraded samples (EOA, MOA, AOA) for direct comparison."
    ])

    # --- Slide 6: Step 4 - Calculation ---
    add_slide("4. Angle Calculation Model", [
        "• Equation: θ = [(H − H₀) / (H₉₀ − H₀)] × 90°",
        "• Method:",
        "    - ROI sliced into 100 normalized depth bins.",
        "    - Circular mean hue calculated for each bin.",
        "    - Linear interpolation determines angle θ.",
        "• Result: Continuous depth-dependent orientation profile."
    ])
    
    # --- Slide 7: Comparison ---
    add_slide("5. Comparative Analysis Framework", [
        "• Groups Analyzed:",
        "    - Healthy Reference (Baseline)",
        "    - Early OA (EOA)",
        "    - Moderate OA (MOA)",
        "    - Advanced OA (AOA)",
        "• Metric: Overlay of Depth vs. Angle curves.",
        "• Goal: Quantify loss of surface tangency and structural disorganization."
    ])

    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    try:
        create_presentation("Cartilage_Methodology_Workflow.pptx")
    except ImportError:
        print("Error: python-pptx library is missing.")
        print("Please run: pip install python-pptx")
    except Exception as e:
        print(f"An error occurred: {e}")
