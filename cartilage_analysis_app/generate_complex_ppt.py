
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_complex_presentation(output_file):
    prs = Presentation()
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    GREEN = RGBColor(0, 128, 0)
    RED = RGBColor(192, 0, 0)
    BLUE = RGBColor(0, 76, 153)
    ORANGE = RGBColor(255, 165, 0)
    GREY = RGBColor(100, 100, 100)
    WHITE = RGBColor(255, 255, 255)

    # --- Slide 1: Comprehensive Workflow ---
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Integrated Methodology: Image Analysis to FE Modeling"
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Helper to create styled box
    def add_box(left, top, width, height, color, title_text, bullets):
        # Background shape (no fill, colored border)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.background() # Transparent
        line = shape.line
        line.color.rgb = color
        line.width = Pt(3)
        
        # Title Background for box
        t_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.4))
        t_bg.fill.solid()
        t_bg.fill.fore_color.rgb = color
        t_bg.line.fill.background()
        t_bg.text_frame.text = title_text
        t_bg.text_frame.paragraphs[0].font.color.rgb = WHITE
        t_bg.text_frame.paragraphs[0].font.bold = True
        t_bg.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Content
        tb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.5), width - Inches(0.2), height - Inches(0.6))
        tf = tb.text_frame
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(10)
            p.level = 0
            p.space_after = Pt(6)

    # === LEFT COLUMN: Input Phase ===
    # PLM Image & ROI
    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(1.5), Inches(1.0)).text_frame.text = "PLM Image Input"
    
    # Arrow Down
    slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.1), Inches(2.6), Inches(0.3), Inches(0.5))
    
    # ROI Segmentation
    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.2), Inches(1.5), Inches(1.0)).text_frame.text = "ROI Selection & Trimming"
    
    # Split Arrows
    # Green Path (Healthy)
    ar1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.1), Inches(3.5), Inches(1.0), Inches(0.3))
    ar1.fill.solid()
    ar1.fill.fore_color.rgb = GREEN
    
    # === MIDDLE COLUMN: Image Processing ===
    
    # Box 1: Healthy Reference (Green)
    add_box(Inches(3.2), Inches(1.5), Inches(4.5), Inches(2.8), GREEN, "Healthy Reference (Control)", [
        "1. HSV Conversion & Filtering",
        "2. Reference Hue Detection:",
        "    - SZ (0° Anchor)",
        "    - DZ (90° Anchor)",
        "3. Calibration Map Creation",
        "4. Standardized Partitioning:"
    ])
    
    # Box 2: Degraded Sample (Red)
    add_box(Inches(3.2), Inches(4.5), Inches(4.5), Inches(2.8), RED, "Degraded Sample (OA)", [
        "1. Apply Calibration Map",
        "2. ROI Segmentation (same boundaries)",
        "3. Fiber Angle Calculation:",
        "    - Hue Interpolation",
        "4. Profile Extraction (Depth vs Angle)"
    ])
    
    # === RIGHT COLUMN: Modeling & Validation ===
    
    # Box 3: FE Model (Orange) - Top Right
    add_box(Inches(8.0), Inches(1.5), Inches(4.8), Inches(3.0), ORANGE, "FE Modeling & RVE", [
        "• RVE Generation per Zone (SZ, MZ, DZ)",
        "• Homogenization of Properties",
        "• 2D Axi-symmetric Model",
        "• Stress/Strain Analysis"
    ])
    
    # Box 4: Analysis (Blue) - Bottom Right
    add_box(Inches(8.0), Inches(4.7), Inches(4.8), Inches(2.6), BLUE, "Quantitative Validation", [
        "• Comparative Overlay Graphs",
        "• Delta Analysis (Control vs Degraded)",
        "• Statistical Significance",
        "• Validation against Literature"
    ])

    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    try:
        create_complex_presentation("Complex_Workflow_Layout.pptx")
    except ImportError:
        print("Error: python-pptx library is missing.")
        print("Please run: pip install python-pptx")
    except Exception as e:
        print(f"An error occurred: {e}")
